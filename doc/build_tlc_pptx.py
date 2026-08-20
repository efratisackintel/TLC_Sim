#!/usr/bin/env python3
"""Build a 'TLC for Dummies' PowerPoint - big visuals, plain-language analogies.

Reuses the PNG diagrams in ./assets (run build_tlc_assets.py first).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "TLC_ForDummies.pptx")

BLUE = RGBColor(0x00, 0x71, 0xC5)
DARK = RGBColor(0x1F, 0x2A, 0x44)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
ORANGE = RGBColor(0xE8, 0x82, 0x0C)
LIGHT = RGBColor(0xEA, 0xF3, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x55, 0x55)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space=4, wrap=True):
    """runs: list of (string, size, color, bold) or a plain string."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [[(runs, 18, DARK, False)]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        if isinstance(para, str):
            para = [(para, 18, DARK, False)]
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = "Segoe UI"
    return tb


def banner(s, title, sub=None):
    rect(s, 0, 0, 13.333, 1.15, BLUE, rounded=False)
    text(s, 0.5, 0.12, 12.3, 0.9,
         [[(title, 30, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(s, 0.5, 0.78, 12.3, 0.4, [[(sub, 14, RGBColor(0xD6, 0xE9, 0xF9), False)]])


def pic(s, name, x, y, w):
    path = os.path.join(ASSETS, name)
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


def bullets(s, x, y, w, h, items, size=18, color=DARK, gap=10):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
            paras.append([("\u2022  ", size, BLUE, True), (lead, size, color, True),
                          (rest, size, color, False)])
        else:
            paras.append([("\u2022  ", size, BLUE, True), (it, size, color, False)])
    text(s, x, y, w, h, paras, space=gap)


def chip(s, x, y, w, h, label, val, color):
    rect(s, x, y, w, h, LIGHT, line=color)
    text(s, x + 0.1, y + 0.12, w - 0.2, 0.5, [[(label, 13, GREY, True)]], align=PP_ALIGN.CENTER)
    text(s, x + 0.1, y + 0.5, w - 0.2, h - 0.6, [[(val, 20, color, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---------------- Slide 1: Title ----------------
s = slide()
rect(s, 0, 0, 13.333, 7.5, DARK, rounded=False)
rect(s, 0, 2.55, 13.333, 2.4, BLUE, rounded=False)
text(s, 0.8, 2.75, 11.7, 1.2, [[("TLC for Dummies", 48, WHITE, True)]],
     anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.8, 3.95, 11.7, 0.8,
     [[("How WiFi firmware picks the fastest reliable speed \u2014 in plain English", 22, RGBColor(0xD6,0xE9,0xF9), False)]])
text(s, 0.8, 6.4, 11.7, 0.6, [[("Transmit Link Control \u2022 Rate Scale Manager \u2022 2026", 16, RGBColor(0x9F,0xB4,0xCC), False)]])

# ---------------- Slide 2: What is TLC (analogy) ----------------
s = slide()
banner(s, "What is TLC?", "Think of a driver choosing a gear")
text(s, 0.5, 1.4, 7.4, 5.6, [
    [("The problem", 22, BLUE, True)],
    [("WiFi can send data at many speeds. Fast speeds need a clean, strong signal. "
      "If the signal is weak, fast data gets lost and must be re-sent \u2014 which is slower.", 18, DARK, False)],
    [("", 8, DARK, False)],
    [("What TLC does", 22, BLUE, True)],
    [("TLC is the \u201cauto gearbox\u201d of WiFi. It constantly watches how many packets "
      "get through and shifts the speed up or down to stay fast AND reliable.", 18, DARK, False)],
    [("", 8, DARK, False)],
    [("The golden rule", 22, GREEN, True)],
    [("Send as fast as possible \u2014 but only as fast as the link can reliably carry.", 18, DARK, False)],
], space=8)
rect(s, 8.3, 1.5, 4.5, 5.2, LIGHT, line=BLUE)
text(s, 8.5, 1.7, 4.1, 5.0, [
    [("\U0001F697  Driving analogy", 20, DARK, True)],
    [("", 6, DARK, False)],
    [("Open road (strong signal)", 16, GREEN, True)],
    [("\u2192 high gear, go fast (high MCS, MIMO)", 15, DARK, False)],
    [("", 6, DARK, False)],
    [("Traffic / hills (weak signal)", 16, ORANGE, True)],
    [("\u2192 low gear, stay safe (low MCS, SISO)", 15, DARK, False)],
    [("", 6, DARK, False)],
    [("Skidding (packets lost)", 16, RGBColor(0xC0,0x39,0x2B), True)],
    [("\u2192 shift down immediately", 15, DARK, False)],
], space=6)

# ---------------- Slide 3: The 4 layers ----------------
s = slide()
banner(s, "Who does what?", "Four layers working together")
pic(s, "hierarchy.png", 0.6, 1.4, 7.2)
text(s, 8.1, 1.5, 4.9, 5.6, [
    [("Driver", 18, BLUE, True), (" \u2014 tells TLC what the device can do", 16, DARK, False)],
    [("", 6, DARK, False)],
    [("Upper MAC (the brain)", 18, BLUE, True), (" \u2014 decides the speed", 16, DARK, False)],
    [("", 6, DARK, False)],
    [("Lower MAC", 18, GREEN, True), (" \u2014 applies the speed & counts wins/losses", 16, DARK, False)],
    [("", 6, DARK, False)],
    [("PHY / Radio", 18, ORANGE, True), (" \u2014 actually transmits", 16, DARK, False)],
    [("", 10, DARK, False)],
    [("Control flows down \u2b07, results flow back up \u2b06.", 16, GREY, True)],
], space=6)

# ---------------- Slide 4: Feedback loop ----------------
s = slide()
banner(s, "The magic: a feedback loop", "Try \u2192 Measure \u2192 Adjust \u2192 repeat")
pic(s, "feedback.png", 1.4, 1.6, 10.5)
text(s, 1.0, 5.7, 11.3, 1.3, [
    [("Every time a batch of packets is sent, TLC checks how many were acknowledged, "
      "recalculates the ", 18, DARK, False),
     ("Success Ratio", 18, GREEN, True),
     (", and re-tunes the next transmission. This repeats thousands of times per second.", 18, DARK, False)],
], space=6)

# ---------------- Slide 5: Inputs ----------------
s = slide()
banner(s, "What TLC watches (Inputs)", "Mostly one number rules them all")
rect(s, 0.6, 1.45, 12.1, 1.9, LIGHT, line=GREEN)
text(s, 0.85, 1.6, 11.6, 1.7, [
    [("\u2b50 Success Ratio (SR)", 22, GREEN, True), ("  = packets acknowledged \u00f7 packets sent", 18, DARK, False)],
    [("Like a mail carrier tracking \u201chow many letters were delivered vs. sent.\u201d "
      "High SR = link is healthy = try faster. Low SR = slow down.", 16, DARK, False)],
], space=6)
data = [
    ("Success Ratio", "0-100%", "how many got through", GREEN),
    ("PER", "= 100% - SR", "how many were lost", ORANGE),
    ("RSSI / SNR", "signal strength", "helps pick a start speed", BLUE),
    ("Timers", "~10 s idle", "re-check even when quiet", DARK),
]
x = 0.6
for label, val, desc, color in data:
    chip(s, x, 3.7, 2.85, 1.5, label, val, color)
    text(s, x, 5.25, 2.85, 0.9, [[(desc, 13, GREY, False)]], align=PP_ALIGN.CENTER)
    x += 3.06
text(s, 0.6, 6.4, 12.1, 0.6,
     [[("Refresh: every ~20 packets normally; only every ~2000 once locked on the best speed (saves power).", 15, GREY, True)]])

# ---------------- Slide 6: Outputs ----------------
s = slide()
banner(s, "What TLC decides (Outputs)", "The recipe for each transmission")
outs = [
    ("Speed (MCS)", "how dense the data is", BLUE),
    ("Streams (NSS)", "1 lane (SISO) or 2 (MIMO)", GREEN),
    ("Bandwidth", "20 - 320 MHz road width", BLUE),
    ("Guard Interval", "gap between symbols", ORANGE),
    ("TX Power", "turn down when close", GREEN),
    ("Aggregation", "bundle packets together", BLUE),
]
x, y = 0.6, 1.55
for i, (label, desc, color) in enumerate(outs):
    cx = 0.6 + (i % 3) * 4.1
    cy = 1.55 + (i // 3) * 1.75
    rect(s, cx, cy, 3.85, 1.55, WHITE, line=color)
    text(s, cx + 0.15, cy + 0.12, 3.55, 0.6, [[(label, 18, color, True)]])
    text(s, cx + 0.15, cy + 0.72, 3.55, 0.75, [[(desc, 15, DARK, False)]])
text(s, 0.6, 5.3, 12.1, 1.6, [
    [("All of these are packed into one number called ", 17, DARK, False),
     ("rate_n_flags", 17, BLUE, True),
     (" that the radio understands, plus a ", 17, DARK, False),
     ("retry table", 17, GREEN, True), (" (plan B, C, D\u2026).", 17, DARK, False)],
], space=6)

# ---------------- Slide 7: Retry ladder ----------------
s = slide()
banner(s, "Plan B, C, D\u2026 the Retry Table", "If the fast speed fails, fall back instantly")
pic(s, "retry.png", 1.6, 1.5, 7.0)
text(s, 8.9, 1.7, 4.0, 5.2, [
    [("Why?", 20, BLUE, True)],
    [("If the best speed fails for a packet, the radio doesn\u2019t give up \u2014 "
      "it instantly retries at a slower, safer speed.", 16, DARK, False)],
    [("", 8, DARK, False)],
    [("Top of ladder", 16, GREEN, True), (" = fastest", 15, DARK, False)],
    [("Bottom of ladder", 16, ORANGE, True), (" = most robust (legacy)", 15, DARK, False)],
    [("", 8, DARK, False)],
    [("This keeps data moving even during a quick fade, before the algorithm reacts.", 16, GREY, False)],
], space=6)

# ---------------- Slide 8: State machine ----------------
s = slide()
banner(s, "TLC's three moods", "Cruise \u2022 Explore \u2022 Save power")
pic(s, "state.png", 1.4, 1.5, 10.5)
cols = [
    ("STAY (cruise)", "Happy on a proven speed. Small nudges up or down.", BLUE),
    ("SEARCH (explore)", "Curious \u2014 tries other speeds/widths to find something better.", GREEN),
    ("TPC (save power)", "At top speed already \u2014 now trims TX power to save energy.", ORANGE),
]
x = 0.6
for label, desc, color in cols:
    rect(s, x, 5.3, 3.9, 1.65, WHITE, line=color)
    text(s, x + 0.15, 5.42, 3.6, 0.55, [[(label, 17, color, True)]])
    text(s, x + 0.15, 5.98, 3.6, 0.9, [[(desc, 14, DARK, False)]])
    x += 4.13

# ---------------- Slide 9: Throughput chart ----------------
s = slide()
banner(s, "Faster settings = more throughput", "Real firmware numbers (HE, 80 MHz)")
pic(s, "tpt_chart.png", 1.5, 1.5, 10.3)
text(s, 1.0, 6.2, 11.3, 1.0, [
    [("Higher MCS and 2 streams (MIMO) roughly double the speed \u2014 but need a cleaner signal. "
      "TLC climbs this chart only as far as the link allows.", 17, DARK, False)],
], space=6)

# ---------------- Slide 10: Cheat sheet ----------------
s = slide()
banner(s, "Cheat sheet: the numbers that matter")
rows = [
    ("95%", "\u201cPerfect\u201d success ratio \u2014 above this, never slow down"),
    ("90% / 15%", "No-decrease / force-decrease thresholds"),
    ("94% / 50%", "Turn A-MSDU bundling on / off"),
    ("200 ms", "Minimum time between speed-up attempts"),
    ("20 \u2192 2000", "Packets between checks (normal \u2192 locked on best)"),
    ("2 dB", "Power back-off step when saving energy (TPC)"),
]
y = 1.5
for i, (num, desc) in enumerate(rows):
    fill = LIGHT if i % 2 == 0 else WHITE
    rect(s, 0.8, y, 11.7, 0.82, fill, line=BLUE)
    text(s, 1.0, y + 0.06, 2.6, 0.7, [[(num, 22, BLUE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 3.7, y + 0.06, 8.6, 0.7, [[(desc, 17, DARK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.9

# ---------------- Slide 11: Try it ----------------
s = slide()
rect(s, 0, 0, 13.333, 7.5, DARK, rounded=False)
rect(s, 0, 2.7, 13.333, 2.1, GREEN, rounded=False)
text(s, 0.8, 2.9, 11.7, 1.1, [[("\U0001F3AE  Now play with it!", 40, WHITE, True)]],
     anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.8, 5.1, 11.7, 1.8, [
    [("Open ", 22, WHITE, False), ("tlc_playground.html", 22, RGBColor(0xBF, 0xFF, 0xBF), True),
     (" in any browser.", 22, WHITE, False)],
    [("Drag the signal-quality slider and watch TLC choose the speed, width, streams, "
      "power and retry ladder \u2014 live.", 18, RGBColor(0xD6,0xF0,0xDE), False)],
], space=10)

prs.save(OUT)
print("Saved:", OUT)
